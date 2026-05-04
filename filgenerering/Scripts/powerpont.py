from pptx import Presentation

def create_presentation():
    while True:
        try:
            # Ber om tittle for presangtasjonen
            prs_title = input("Hva skal presentasjonen hete? ").strip().lower().replace(" ", "_")
            if prs_title.isdigit():
                print("Vennligst skriv inn en gyldig tittel for presentasjonen.")
                continue
            # Ber om antall slides i presentasjonen
            num_slides = int(input("Hvor mange slides vil du ha i presentasjonen? "))
            print(f"Du har valgt å lage en presentasjon med {num_slides} slides.")

            # Ber om tittel og innhold for hver slide
            slide_content = []
            for i in range(num_slides):
                title = input(f"Skriv inn tittelen for slide {i + 1}: ")
                content = input(f"Skriv inn innholdet for slide {i + 1}: ")
                slide_content.append((title, content))  
            break      
        except ValueError:
            print("Vennligst skriv inn et gyldig tall.")

    

    # Oppretter en ny presentasjon
    prs = Presentation()
    
    # Løkke for å hente antall slides og legge til tittel og innhold for hver slide
    for i in range(num_slides):
        slide_layout = prs.slide_layouts[0] # Velger en enkel layout for hver slide
        slide = prs.slides.add_slide(slide_layout) # Legger til en ny slide i presentasjonen
        title = slide.shapes.title # Får tilgang til tittelen på sliden
        title.text = slide_content[i][0]  # Setter tittelen for sliden
        content = slide.shapes.placeholders[1] # Får tilgang til innholdet på sliden
        content.text = slide_content[i][1]  # Setter innholdet for sliden
    prs.save(f"{prs_title}.pptx") # Lagerer presentasjonen med det angitte navnet
    print(f"Presentasjonen er lagret som '{prs_title}.pptx'.")
    return f"{prs_title}.pptx"


create_presentation()